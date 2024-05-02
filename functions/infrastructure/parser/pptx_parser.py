from application.parser import IParser, ParsingResult
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from application.parser.iparser import DocSection
from typing import List, Any

class PPTXParser(IParser): 
    ...
    def parse(self, file: BytesIO) -> ParsingResult:
        text_sections: List[DocSection[str]] = []
        image_sections: List[DocSection[Any]] = []

        prs = Presentation(file)

                # Current index to track order
        current_index = 0

        # Loop through slides and extract text and images
        for slide in prs.slides:
            for shape in slide.shapes:
                # Handle text sections
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text:
                            text_sections.append(DocSection[str](index=current_index, content=paragraph.text))
                            current_index += 1
                
                # Handle image sections
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = BytesIO(shape.image.blob)
                    image_sections.append(DocSection[Any](index=current_index, content=image_bytes))
                    current_index += 1
                
                # Handle group shapes with nested images
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for s in shape.shapes:
                        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_bytes = BytesIO(s.image.blob)
                            image_sections.append(DocSection[Any](index=current_index, content=image_bytes))
                            current_index += 1  

        all_sections = sorted(text_sections + image_sections, key=lambda ds: ds.index)
        print("hola")
        # Print all sections in order
        for section in all_sections:
            if isinstance(section.content, str):
                print(f"Text {section.index}: {section.content}")
            elif isinstance(section.content, BytesIO):
                print(f"Image {section.index}")
                

        return ParsingResult(text_sections=text_sections,image_sections=image_sections)